#!/usr/bin/env python3
from dotenv import load_dotenv
load_dotenv()

import anthropic
import sys
import os

client = anthropic.Anthropic()

system_prompt = """
You are an experienced Senior AI Product Manager and strategic thought partner working directly with Kurren.

About Kurren:
- Senior AI PM at Entain, a global sports betting and entertainment company
- Owns personalisation and AI across all product areas (Sportsbook, Gaming, Risk, Commercial, Internal Platforms)
- Role is roughly 50/50 enablement and product delivery
- Reports to the Head of AI
- Background in consumer marketplace product (Depop, iOS discovery)
- Strong in experimentation, personalisation, cross-functional leadership
- Weaker in: sports betting domain, risk as a product area, internal platform thinking

How Kurren works best:
- Needs help structuring thinking, not just validating it
- Prefers direct, concise communication — no fluff
- Responds well to being challenged and pushed back on
- Thinks in terms of evidence, hypotheses, and measurable outcomes

Your primary jobs:
1. Help Kurren structure messy thinking into clear strategies and plans
2. Draft and critique PRDs, stakeholder updates, roadmaps, and strategy docs
3. Challenge weak assumptions before they become bad decisions
4. Help Kurren build credibility and domain knowledge in sports betting and AI

Document style guidelines:
- PRDs: clear problem statement, hypothesis, success metrics, scope, risks
- Stakeholder updates: short, outcome-focused, no jargon
- Roadmaps: tied to business outcomes, not feature lists
- Strategy docs: diagnosis first, then approach, then execution

Entain domain context:
- Key metrics: GGR, NGR, handle, hold %, player lifetime value
- Regulatory environment: UKGC, responsible gambling, affordability checks
- Key tension: commercial growth vs. risk and compliance
- AI use cases: personalisation, risk/fraud, customer support, internal tooling

Always:
- Ask clarifying questions before diving into a task if the brief is unclear
- Flag assumptions you're making
- End substantive responses with a clear "next step" or question to keep momentum
"""

conversation_history = []

# Auto-load context file if it exists
context_file = os.path.join(os.path.dirname(__file__), 'context.md')
if os.path.exists(context_file):
    with open(context_file, 'r') as f:
        context_contents = f.read()
    conversation_history.append({
        "role": "user",
        "content": f"Here is my current context for this session:\n\n{context_contents}"
    })
    conversation_history.append({
        "role": "assistant",
        "content": "Got it — I have your current context. What do you want to work on?"
    })

# Check if a file was passed in
def read_file(filename):
    if filename.endswith('.txt') or filename.endswith('.md'):
        with open(filename, 'r') as f:
            return f.read()
    elif filename.endswith('.pdf'):
        from pypdf import PdfReader
        reader = PdfReader(filename)
        return "\n".join([page.extract_text() for page in reader.pages])
    elif filename.endswith('.docx'):
        from docx import Document
        doc = Document(filename)
        return "\n".join([para.text for para in doc.paragraphs])
    elif filename.endswith('.pptx'):
        from pptx import Presentation
        prs = Presentation(filename)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    else:
        print(f"Unsupported file type: {filename}")
        sys.exit(1)

if len(sys.argv) > 1:
    all_contents = []
    for filename in sys.argv[1:]:
        try:
            file_contents = read_file(filename)
            print(f"📄 Loaded: {filename}")
            all_contents.append(f"--- {filename} ---\n{file_contents}")
        except FileNotFoundError:
            print(f"File not found: {filename}")
            sys.exit(1)

    conversation_history.append({
        "role": "user",
        "content": f"I'm sharing these documents with you:\n\n" + "\n\n".join(all_contents)
    })

print("AI PM Assistant ready. Type 'quit' to exit.")
print("---")

while True:
    user_input = input("You: ")

    if user_input.lower() == "quit":
        print("Ending conversation.")
        break

    if user_input.lower() == "save":
        print("Updating context file...")
        summary_request = [{
            "role": "user",
            "content": """Please update my context.md file based on our conversation. 
Keep the same format but update:
- Current priorities
- Open questions
- Key decisions made
- This week's focus

Return ONLY the markdown content, nothing else."""
        }]
        summary = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=1024,
            system=system_prompt,
            messages=conversation_history + summary_request
        )
        with open(context_file, 'w') as f:
            f.write(summary.content[0].text)
        print("✅ Context file updated. Type 'quit' to exit or keep chatting.")
        continue

    if not user_input.strip():
        continue

    conversation_history.append({
        "role": "user",
        "content": user_input
    })

    message = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=4096,
        system=system_prompt,
        messages=conversation_history
    )

    response = message.content[0].text

    conversation_history.append({
        "role": "assistant",
        "content": response
    })

    print(f"Claude: {response}")
    print("---")

    # Save response to file
    with open("responses.md", "a") as f:
        f.write(f"## You:\n{user_input}\n\n## Claude:\n{response}\n\n---\n\n")